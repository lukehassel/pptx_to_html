from pptx import Presentation
from pptx.shapes.picture import Picture

from src.pptx_to_html.data.datasources.generate_html import generate_html
from src.pptx_to_html.domain.entities.elements.Image import Image
from src.pptx_to_html.domain.entities.elements.Text import Text
from src.pptx_to_html.domain.entities.elements.slide import Slide


class PowerPointInterface:

    def get_content(self, path) -> str:
        pass


class PowerPoint(PowerPointInterface):

    def get_content(self, path) -> str:
        text = self.read(path)
        html = generate_html(text)

        return html

    def read(self, path):
        prs = Presentation(path)

        text_runs = []
        content = []

        for slide in prs.slides:

            shapes = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shapes.append(shape)
                if isinstance(shape, Picture):
                    shapes.append(shape)
                # sorted_list =sorted(shape.text_frame.paragraphs, )
            if len(shapes) > 0:
                ordered_shapes = sorted(shapes, key=lambda x: x.top)

                content.append(Slide())
                for shape in ordered_shapes:

                    # Text
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if len(paragraph.text) > 0:
                                text_runs.append(paragraph.text)
                                average_size = self.get_average_size(paragraph, shape)
                                # print(average_size, ":", paragraph.text)

                                t = Text(text=paragraph.text, size=average_size)
                                content.append(t)
                    # pictures
                    if isinstance(shape, Picture):
                        # print(shape.image.blob)<
                        t = Image(bytes=shape.image.blob)
                        content.append(t)

        return content

    def get_average_size(self, paragraph, shape):
        count = 0
        adder = 0
        if paragraph.font.size is None:
            for run in paragraph.runs:
                if not run.font is None and not run.font.size is None:
                    count += 1
                    adder += run.font.size.pt
            if count > 0 and adder > 0:
                return adder / count
            else:
                # print(shape.height.pt/2, " :", paragraph.text.count('\n')," " ,paragraph.text)
                return shape.height.pt / 2
        else:
            return paragraph.font.size
